import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def create_deck():
    prs = Presentation()
    # Use widescreen 16:9 layout
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Color Palette: Bright Light Theme
    # Background: White / Very Light Gray (250, 250, 252)
    # Primary (Dark Slate Navy): 24, 43, 73
    # Accent (Teal Blue): 0, 150, 167
    # Text (Dark Charcoal): 40, 40, 40
    # Secondary Accent (Sky Blue): 3, 169, 244
    
    C_BG = RGBColor(250, 250, 252)
    C_PRIMARY = RGBColor(24, 43, 73)
    C_ACCENT = RGBColor(0, 150, 167)
    C_TEXT = RGBColor(40, 40, 40)
    C_SECONDARY = RGBColor(3, 169, 244)
    C_WHITE = RGBColor(255, 255, 255)
    C_MUTED = RGBColor(120, 120, 120)

    # Helper function to apply solid background color to slides
    def set_slide_background(slide, color):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color

    # Helper to add standard slide title
    def add_slide_header(slide, title_text):
        title_box = slide.shapes.add_textbox(Inches(0.75), Inches(0.4), Inches(11.83), Inches(0.9))
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.name = "Segoe UI"
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = C_PRIMARY

        # Add a subtle accent bar under the title
        accent_bar = slide.shapes.add_shape(
            1, # MSO_SHAPE.RECTANGLE
            Inches(0.75), Inches(1.3), Inches(1.5), Inches(0.06)
        )
        accent_bar.fill.solid()
        accent_bar.fill.fore_color.rgb = C_ACCENT
        accent_bar.line.color.rgb = C_ACCENT

    # Helper to format paragraph
    def format_p(p, text, font_size, bold=False, color=C_TEXT, italic=False, space_after=12):
        p.text = text
        p.font.name = "Segoe UI"
        p.font.size = Pt(font_size)
        p.font.bold = bold
        p.font.italic = italic
        p.font.color.rgb = color
        p.space_after = Pt(space_after)

    # ==================== SLIDE 1: TITLE SLIDE ====================
    slide_layout = prs.slide_layouts[6] # Blank layout
    slide1 = prs.slides.add_slide(slide_layout)
    set_slide_background(slide1, C_BG)

    # Left colored accent panel
    left_panel = slide1.shapes.add_shape(
        1, Inches(0), Inches(0), Inches(0.4), Inches(7.5)
    )
    left_panel.fill.solid()
    left_panel.fill.fore_color.rgb = C_ACCENT
    left_panel.line.color.rgb = C_ACCENT

    # Main Title Box
    title_box = slide1.shapes.add_textbox(Inches(1.0), Inches(0.8), Inches(11.3), Inches(1.8))
    tf1 = title_box.text_frame
    tf1.word_wrap = True
    p1 = tf1.paragraphs[0]
    format_p(p1, "AI-Driven Fake Job Posting Detection System", 44, bold=True, color=C_PRIMARY)
    p2 = tf1.add_paragraph()
    format_p(p2, "Classifying Fraudulent Job Advertisements using ML & NLP", 20, bold=False, color=C_ACCENT)

    # Submitted By Details (Left column)
    submit_box = slide1.shapes.add_textbox(Inches(1.0), Inches(3.0), Inches(5.5), Inches(3.5))
    tf_sub = submit_box.text_frame
    tf_sub.word_wrap = True
    format_p(tf_sub.paragraphs[0], "SUBMITTED BY:", 12, bold=True, color=C_MUTED, space_after=6)
    
    names = [
        ("Zeel Dharmendrakumar Parmar", "70438497"),
        ("Yash Sanjaybhai Patel", "38402834"),
        ("Bhavani manjunatha", "68681898"),
        ("Ashik Kirmani", "78221471")
    ]
    for name, student_id in names:
        p = tf_sub.add_paragraph()
        format_p(p, f"•  {name} ({student_id})", 14, bold=True, color=C_PRIMARY, space_after=4)

    p_course = tf_sub.add_paragraph()
    p_course.space_before = Pt(12)
    format_p(p_course, "MSc Software Engineering", 14, bold=False, color=C_TEXT)

    # University & Submission Details (Right column)
    univ_box = slide1.shapes.add_textbox(Inches(7.0), Inches(3.0), Inches(5.3), Inches(3.5))
    tf_univ = univ_box.text_frame
    tf_univ.word_wrap = True
    
    format_p(tf_univ.paragraphs[0], "UNIVERSITY:", 12, bold=True, color=C_MUTED, space_after=6)
    p_u1 = tf_univ.add_paragraph()
    format_p(p_u1, "University of Europe for Applied Sciences", 16, bold=True, color=C_PRIMARY, space_after=2)
    p_u2 = tf_univ.add_paragraph()
    format_p(p_u2, "Potsdam, Germany", 14, bold=False, color=C_TEXT, space_after=18)

    p_subto = tf_univ.add_paragraph()
    format_p(p_subto, "SUBMITTED TO:", 12, bold=True, color=C_MUTED, space_after=6)
    p_prof = tf_univ.add_paragraph()
    format_p(p_prof, "Prof. Raza Ali", 16, bold=True, color=C_PRIMARY, space_after=2)
    p_email = tf_univ.add_paragraph()
    format_p(p_email, "raza.ali@ue-germany.de", 13, bold=False, color=C_ACCENT)

    # ==================== SLIDE 2: PROBLEM SETTING ====================
    slide2 = prs.slides.add_slide(slide_layout)
    set_slide_background(slide2, C_BG)
    add_slide_header(slide2, "Problem Setting & Motivation")

    # Content split into left visual card and right details card
    # Left Box
    left_box = slide2.shapes.add_shape(1, Inches(0.75), Inches(1.8), Inches(4.5), Inches(4.8))
    left_box.fill.solid()
    left_box.fill.fore_color.rgb = C_PRIMARY
    left_box.line.fill.background()
    
    tf_left = left_box.text_frame
    tf_left.word_wrap = True
    tf_left.margin_left = Inches(0.4)
    tf_left.margin_right = Inches(0.4)
    tf_left.margin_top = Inches(0.5)
    
    p = tf_left.paragraphs[0]
    format_p(p, "THE THREAT", 12, bold=True, color=C_SECONDARY, space_after=12)
    p_body = tf_left.add_paragraph()
    format_p(p_body, "Online job boards host millions of listings. Scammers exploit these platforms to post fake jobs for:", 14, bold=False, color=C_WHITE, space_after=18)
    
    bullet_items = [
        "Identity Theft (personal data)",
        "Financial Scams (advance fees)",
        "Credential Harvesting"
    ]
    for item in bullet_items:
        p_item = tf_left.add_paragraph()
        format_p(p_item, f"•  {item}", 14, bold=True, color=C_WHITE, space_after=8)

    # Right Content Box
    right_box = slide2.shapes.add_textbox(Inches(5.6), Inches(1.8), Inches(7.0), Inches(4.8))
    tf_right = right_box.text_frame
    tf_right.word_wrap = True
    
    format_p(tf_right.paragraphs[0], "OUR CORE OBJECTIVES", 14, bold=True, color=C_ACCENT, space_after=12)
    
    obj_items = [
        ("Binary Classification Task", "Predict whether a posting is Legitimate (0) or Fraudulent (1)."),
        ("Dual-Feature Fusion", "Analyze both raw text details (NLP) and listing metadata characteristics."),
        ("Explainable Flagging", "Protect job seekers and reduce manual review costs for recruitment platforms.")
    ]
    for title, desc in obj_items:
        p_t = tf_right.add_paragraph()
        format_p(p_t, f"✔  {title}", 15, bold=True, color=C_PRIMARY, space_after=2)
        p_d = tf_right.add_paragraph()
        format_p(p_d, desc, 13, bold=False, color=C_TEXT, space_after=12)

    # ==================== SLIDE 3: DATASET OVERVIEW ====================
    slide3 = prs.slides.add_slide(slide_layout)
    set_slide_background(slide3, C_BG)
    add_slide_header(slide3, "Dataset & Class Imbalance Challenge")

    # Left Stats Card
    card_left = slide3.shapes.add_shape(1, Inches(0.75), Inches(1.8), Inches(5.2), Inches(4.8))
    card_left.fill.solid()
    card_left.fill.fore_color.rgb = C_WHITE
    card_left.line.color.rgb = RGBColor(220, 220, 225)
    
    tf_card = card_left.text_frame
    tf_card.word_wrap = True
    tf_card.margin_left = Inches(0.4)
    tf_card.margin_top = Inches(0.4)
    
    format_p(tf_card.paragraphs[0], "DATASET PROFILE", 13, bold=True, color=C_ACCENT, space_after=14)
    
    profile_items = [
        ("Source", "Employment Scam Aegean Dataset (EMSCAD)"),
        ("Dimensions", "17,880 total job postings"),
        ("Text Features", "Job Title, Description, Profile, Benefits, Requirements"),
        ("Metadata", "Telecommuting, Logo, Questions, Experience, Education")
    ]
    for label, val in profile_items:
        p_l = tf_card.add_paragraph()
        format_p(p_l, f"{label}:", 11, bold=True, color=C_MUTED, space_after=2)
        p_v = tf_card.add_paragraph()
        format_p(p_v, val, 14, bold=True, color=C_PRIMARY, space_after=8)

    # Right Class Imbalance Card
    card_right = slide3.shapes.add_textbox(Inches(6.4), Inches(1.8), Inches(6.2), Inches(4.8))
    tf_cr = card_right.text_frame
    tf_cr.word_wrap = True
    
    format_p(tf_cr.paragraphs[0], "THE CLASS IMBALANCE TRAP", 14, bold=True, color=C_PRIMARY, space_after=12)
    
    p_desc = tf_cr.add_paragraph()
    format_p(p_desc, "Legitimate jobs represent 95.16% (17,014 records), while Fraudulent postings are only 4.84% (866 records).", 14, bold=False, color=C_TEXT, space_after=14)
    
    # Large numbers callout
    p_call = tf_cr.add_paragraph()
    format_p(p_call, "95.16% vs 4.84%", 26, bold=True, color=C_ACCENT, space_after=12)
    
    p_conseq = tf_cr.add_paragraph()
    format_p(p_conseq, "Consequences for Evaluation:", 13, bold=True, color=C_PRIMARY, space_after=4)
    p_conseq_d = tf_cr.add_paragraph()
    format_p(p_conseq_d, "Traditional accuracy is misleading. If a model predicts 'Legitimate' for all posts, it achieves 95.16% accuracy but catches 0 scams. F1-Score, Precision, and Recall must be used instead.", 13, bold=False, color=C_TEXT)

    # ==================== SLIDE 4: NLP PREPROCESSING ====================
    slide4 = prs.slides.add_slide(slide_layout)
    set_slide_background(slide4, C_BG)
    add_slide_header(slide4, "Natural Language Processing (NLP) Pipeline")

    # 4 distinct steps arranged horizontally as cards
    card_w = Inches(2.7)
    card_h = Inches(4.5)
    gap = Inches(0.3)
    start_x = Inches(0.75)
    start_y = Inches(2.0)
    
    steps = [
        ("1. Aggregation", "Concatenated text fields (Title, Description, Profile, Benefits, Requirements) into a single text block for each listing."),
        ("2. Cleaning & Lowering", "Converted text to lowercase and removed special characters, punctuation, and numerical digits using Regex patterns."),
        ("3. Tokenization & Stopwords", "Split paragraphs into word tokens and filtered out common English stopwords (e.g., 'the', 'is', 'and') that lack classification weight."),
        ("4. Lemmatization", "Reduced word tokens to their root dictionary form (e.g., 'requirements' -> 'requirement', 'benefits' -> 'benefit') using NLTK WordNet.")
    ]
    
    for i, (title, desc) in enumerate(steps):
        x_pos = start_x + i * (card_w + gap)
        card = slide4.shapes.add_shape(1, x_pos, start_y, card_w, card_h)
        card.fill.solid()
        card.fill.fore_color.rgb = C_WHITE
        card.line.color.rgb = RGBColor(220, 220, 225)
        
        tf_step = card.text_frame
        tf_step.word_wrap = True
        tf_step.margin_left = Inches(0.2)
        tf_step.margin_right = Inches(0.2)
        tf_step.margin_top = Inches(0.3)
        
        format_p(tf_step.paragraphs[0], title, 15, bold=True, color=C_PRIMARY, space_after=12)
        p_desc = tf_step.add_paragraph()
        format_p(p_desc, desc, 12, bold=False, color=C_TEXT, space_after=0)

    # ==================== SLIDE 5: FEATURE ENGINEERING ====================
    slide5 = prs.slides.add_slide(slide_layout)
    set_slide_background(slide5, C_BG)
    add_slide_header(slide5, "Feature Engineering & Fusion")

    # Left: TF-IDF Text Features
    f_left = slide5.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(5.5), Inches(4.8))
    tf_f_left = f_left.text_frame
    tf_f_left.word_wrap = True
    format_p(tf_f_left.paragraphs[0], "TEXT REPRESENTATION (TF-IDF)", 16, bold=True, color=C_ACCENT, space_after=12)
    
    tf_items = [
        ("TF-IDF Vectorization", "Translates raw words into numeric vectors based on frequency and cross-document rarity."),
        ("N-gram Range (1, 2)", "Captures single words and two-word sequences (e.g., 'high salary', 'work from home') for contextual indicators."),
        ("Feature Limitation", "Retained the top 5,000 text features to prevent overfitting and optimize dimensionality.")
    ]
    for title, desc in tf_items:
        p_t = tf_f_left.add_paragraph()
        format_p(p_t, f"•  {title}", 14, bold=True, color=C_PRIMARY, space_after=2)
        p_d = tf_f_left.add_paragraph()
        format_p(p_d, desc, 12, bold=False, color=C_TEXT, space_after=10)

    # Right: Metadata & Fusion
    f_right = slide5.shapes.add_textbox(Inches(6.8), Inches(1.8), Inches(5.7), Inches(4.8))
    tf_f_right = f_right.text_frame
    tf_f_right.word_wrap = True
    format_p(tf_f_right.paragraphs[0], "METADATA & FEATURE FUSION", 16, bold=True, color=C_PRIMARY, space_after=12)
    
    fuse_items = [
        ("One-Hot Encoding", "Imputed null categories as 'Unspecified' and one-hot encoded metadata (e.g. required education and experience)."),
        ("Sparse Matrix Stacking", "Fused text TF-IDF vectors and encoded metadata arrays horizontally using scipy.sparse.hstack."),
        ("Why Fusion Works", "Text gives semantic context, while metadata flags structural patterns (e.g., missing company logos).")
    ]
    for title, desc in fuse_items:
        p_t = tf_f_right.add_paragraph()
        format_p(p_t, f"✔  {title}", 14, bold=True, color=C_ACCENT, space_after=2)
        p_d = tf_f_right.add_paragraph()
        format_p(p_d, desc, 12, bold=False, color=C_TEXT, space_after=10)

    # ==================== SLIDE 6: MODELS CHOSEN ====================
    slide6 = prs.slides.add_slide(slide_layout)
    set_slide_background(slide6, C_BG)
    add_slide_header(slide6, "Models & Classification Paradigms")

    # 3 horizontal columns for the 3 models
    col_w = Inches(3.6)
    col_gap = Inches(0.4)
    start_x = Inches(0.75)
    
    models = [
        ("Logistic Regression", "BASELINE (TEXT ONLY)", "Fast, highly interpretable linear classifier. Set a baseline using TF-IDF text features only. Balanced class weights applied.", C_MUTED),
        ("Random Forest", "BAGGING ENSEMBLE", "Generates an ensemble of decision trees. Captures non-linear feature interactions from fused text & metadata. Balanced class weights applied.", C_ACCENT),
        ("XGBoost Classifier", "GRADIENT BOOSTING", "Iterative boosting framework. Leverages scale_pos_weight to assign heavy penalties to minority class (fraudulent) errors.", C_PRIMARY)
    ]
    
    for i, (name, role, desc, color) in enumerate(models):
        x = start_x + i * (col_w + col_gap)
        card = slide6.shapes.add_shape(1, x, Inches(2.0), col_w, Inches(4.5))
        card.fill.solid()
        card.fill.fore_color.rgb = C_WHITE
        card.line.color.rgb = RGBColor(220, 220, 225)
        
        tf_m = card.text_frame
        tf_m.word_wrap = True
        tf_m.margin_left = Inches(0.25)
        tf_m.margin_right = Inches(0.25)
        tf_m.margin_top = Inches(0.4)
        
        format_p(tf_m.paragraphs[0], name, 20, bold=True, color=color, space_after=4)
        p_role = tf_m.add_paragraph()
        format_p(p_role, role, 11, bold=True, color=C_MUTED, space_after=14)
        p_desc = tf_m.add_paragraph()
        format_p(p_desc, desc, 13, bold=False, color=C_TEXT, space_after=0)

    # ==================== SLIDE 7: EVALUATION STRATEGY ====================
    slide7 = prs.slides.add_slide(slide_layout)
    set_slide_background(slide7, C_BG)
    add_slide_header(slide7, "Evaluation Strategy for Imbalanced Data")

    # Left Column
    e_left = slide7.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(5.8), Inches(4.8))
    tf_e_left = e_left.text_frame
    tf_e_left.word_wrap = True
    format_p(tf_e_left.paragraphs[0], "OUR PRIMARY METRICS", 16, bold=True, color=C_PRIMARY, space_after=14)
    
    metrics = [
        ("Precision", "How many flagged jobs are actually fake? (Keeps false alarms low)"),
        ("Recall", "How many of the actual fake jobs did we catch? (Keeps job seekers safe)"),
        ("F1-Score", "Harmonic mean of Precision and Recall. The main metric used to compare models."),
        ("ROC-AUC", "Overall capability to distinguish classes across all probability thresholds.")
    ]
    for title, desc in metrics:
        p_t = tf_e_left.add_paragraph()
        format_p(p_t, f"✔  {title}", 14, bold=True, color=C_ACCENT, space_after=2)
        p_d = tf_e_left.add_paragraph()
        format_p(p_d, desc, 12, bold=False, color=C_TEXT, space_after=8)

    # Right Column
    e_right = slide7.shapes.add_shape(1, Inches(6.8), Inches(1.8), Inches(5.7), Inches(4.8))
    e_right.fill.solid()
    e_right.fill.fore_color.rgb = C_PRIMARY
    e_right.line.fill.background()
    tf_e_right = e_right.text_frame
    tf_e_right.word_wrap = True
    tf_e_right.margin_left = Inches(0.4)
    tf_e_right.margin_top = Inches(0.4)
    tf_e_right.margin_right = Inches(0.4)
    
    format_p(tf_e_right.paragraphs[0], "CLASS WEIGHTING EXPLAINED", 13, bold=True, color=C_SECONDARY, space_after=14)
    p_b1 = tf_e_right.add_paragraph()
    format_p(p_b1, "To prevent models from prioritizing the majority class, we integrated class penalties during training:", 14, bold=False, color=C_WHITE, space_after=14)
    
    p_l1 = tf_e_right.add_paragraph()
    format_p(p_l1, "•  Logistic Regression & Random Forest:", 13, bold=True, color=C_WHITE, space_after=2)
    p_d1 = tf_e_right.add_paragraph()
    format_p(p_d1, "Used class_weight='balanced' to adjust weights inversely proportional to class frequencies.", 12, bold=False, color=C_WHITE, space_after=10)
    
    p_l2 = tf_e_right.add_paragraph()
    format_p(p_l2, "•  XGBoost Classifier:", 13, bold=True, color=C_WHITE, space_after=2)
    p_d2 = tf_e_right.add_paragraph()
    format_p(p_d2, "Calculated scale_pos_weight (~19.64) to multiply the loss of the positive (fake) class errors.", 12, bold=False, color=C_WHITE)

    # ==================== SLIDE 8: RESULTS & COMPARISON ====================
    slide8 = prs.slides.add_slide(slide_layout)
    set_slide_background(slide8, C_BG)
    add_slide_header(slide8, "Model Performance Comparison")

    # Table layout
    # Left Box
    rows = 4
    cols = 6
    left = Inches(0.75)
    top = Inches(1.8)
    width = Inches(11.83)
    height = Inches(2.2)

    table_shape = slide8.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table

    # Column widths
    table.columns[0].width = Inches(3.83) # Model
    for c in range(1, 6):
        table.columns[c].width = Inches(1.6)

    # Headers
    headers = ["Model", "Accuracy", "Precision", "Recall", "F1-Score", "ROC-AUC"]
    for col_idx, header in enumerate(headers):
        cell = table.cell(0, col_idx)
        cell.fill.solid()
        cell.fill.fore_color.rgb = C_PRIMARY
        p = cell.text_frame.paragraphs[0]
        p.text = header
        p.alignment = PP_ALIGN.CENTER if col_idx > 0 else PP_ALIGN.LEFT
        p.font.name = "Segoe UI"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = C_WHITE

    # Table content
    data = [
        ["Logistic Regression (Text Only)", "96.31%", "57.92%", "86.71%", "69.44%", "98.15%"],
        ["Random Forest (Text + Meta)", "98.21%", "92.25%", "68.79%", "78.81%", "99.28%"],
        ["XGBoost Classifier (Text + Meta)", "98.63%", "89.24%", "81.50%", "85.20%", "98.92%"]
    ]

    for row_idx, row_data in enumerate(data):
        for col_idx, val in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.fill.solid()
            # Highlight best F1 model
            if row_idx == 2:
                cell.fill.fore_color.rgb = RGBColor(240, 248, 250) # Light teal
            else:
                cell.fill.fore_color.rgb = C_WHITE
            
            p = cell.text_frame.paragraphs[0]
            p.text = val
            p.alignment = PP_ALIGN.CENTER if col_idx > 0 else PP_ALIGN.LEFT
            p.font.name = "Segoe UI"
            p.font.size = Pt(13)
            p.font.bold = True if row_idx == 2 or col_idx == 0 else False
            p.font.color.rgb = C_ACCENT if row_idx == 2 and col_idx > 0 else C_TEXT

    # Key takeaway box below table
    takeaway = slide8.shapes.add_textbox(Inches(0.75), Inches(4.5), Inches(11.83), Inches(2.2))
    tf_t = takeaway.text_frame
    tf_t.word_wrap = True
    format_p(tf_t.paragraphs[0], "KEY PERFORMANCE FINDINGS", 14, bold=True, color=C_PRIMARY, space_after=8)
    
    p_take1 = tf_t.add_paragraph()
    format_p(p_take1, "•  XGBoost Classifier wins: Achieves the highest F1-Score of 85.20% and robust Recall of 81.50% by combining text context and metadata.", 13, bold=False, color=C_TEXT, space_after=6)
    
    p_take2 = tf_t.add_paragraph()
    format_p(p_take2, "•  Precision vs Recall Tradeoff: Random Forest achieves the highest Precision (92.25%) but misses more scams (Recall 68.79%). XGBoost balances this tradeoff optimally.", 13, bold=False, color=C_TEXT)

    # ==================== SLIDE 9: FEATURE IMPORTANCE ====================
    slide9 = prs.slides.add_slide(slide_layout)
    set_slide_background(slide9, C_BG)
    add_slide_header(slide9, "Feature Importance & Model Interpretability")

    # Left content box
    fi_left = slide9.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(6.0), Inches(4.8))
    tf_fi_l = fi_left.text_frame
    tf_fi_l.word_wrap = True
    format_p(tf_fi_l.paragraphs[0], "EXPLAINABLE FRAUD DETECTORS", 16, bold=True, color=C_PRIMARY, space_after=14)
    
    fi_insights = [
        ("has_company_logo (Highly Predictive)", "Postings without a company logo are statistically much more likely to be fake. Bots or low-effort scam creators often skip uploading logos."),
        ("has_questions", "Fake job postings typically do not include custom screening questions (has_questions = 0) to minimize entry barriers for victims."),
        ("Text Keywords (TF-IDF)", "Certain phrasing such as 'earn money', 'urgently hiring', and 'work from home' paired with unspecified requirements are high-probability scam indicators.")
    ]
    for title, desc in fi_insights:
        p_t = tf_fi_l.add_paragraph()
        format_p(p_t, f"•  {title}", 14, bold=True, color=C_ACCENT, space_after=2)
        p_d = tf_fi_l.add_paragraph()
        format_p(p_d, desc, 12, bold=False, color=C_TEXT, space_after=10)

    # Right visual block
    fi_right = slide9.shapes.add_shape(1, Inches(7.2), Inches(1.8), Inches(5.3), Inches(4.8))
    fi_right.fill.solid()
    fi_right.fill.fore_color.rgb = C_PRIMARY
    fi_right.line.fill.background()
    tf_fi_r = fi_right.text_frame
    tf_fi_r.word_wrap = True
    tf_fi_r.margin_left = Inches(0.4)
    tf_fi_r.margin_top = Inches(0.4)
    tf_fi_r.margin_right = Inches(0.4)
    
    format_p(tf_fi_r.paragraphs[0], "TOP PREDICTIVE FEATURES", 14, bold=True, color=C_SECONDARY, space_after=12)
    
    top_features = [
        "1. has_company_logo (Binary)",
        "2. has_questions (Binary)",
        "3. tfidf_vector_earn_money (Text)",
        "4. tfidf_vector_work_home (Text)",
        "5. required_education_Unspecified (Category)",
        "6. required_experience_Unspecified (Category)"
    ]
    for feat in top_features:
        p_feat = tf_fi_r.add_paragraph()
        format_p(p_feat, feat, 13, bold=True, color=C_WHITE, space_after=8)

    # ==================== SLIDE 10: CONCLUSION ====================
    slide10 = prs.slides.add_slide(slide_layout)
    set_slide_background(slide10, C_BG)
    add_slide_header(slide10, "Conclusion & Future Work")

    # Left: Conclusion
    c_left = slide10.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(5.6), Inches(4.8))
    tf_c_l = c_left.text_frame
    tf_c_l.word_wrap = True
    format_p(tf_c_l.paragraphs[0], "PROJECT SUMMARY", 16, bold=True, color=C_PRIMARY, space_after=14)
    
    conclusions = [
        ("High Detection Capability", "Fusing NLP (TF-IDF) text features with tabular metadata enables the XGBoost classifier to identify over 81.5% of scams."),
        ("Importance of Metadata", "Metadata (logos, questions) acts as a powerful indicator of structural legitimacy, enhancing text-only models.")
    ]
    for title, desc in conclusions:
        p_t = tf_c_l.add_paragraph()
        format_p(p_t, f"✔  {title}", 14, bold=True, color=C_ACCENT, space_after=2)
        p_d = tf_c_l.add_paragraph()
        format_p(p_d, desc, 12, bold=False, color=C_TEXT, space_after=10)

    # Right: Limitations & Future Work
    c_right = slide10.shapes.add_textbox(Inches(6.8), Inches(1.8), Inches(5.7), Inches(4.8))
    tf_c_r = c_right.text_frame
    tf_c_r.word_wrap = True
    format_p(tf_c_r.paragraphs[0], "LIMITATIONS & FUTURE WORK", 16, bold=True, color=C_PRIMARY, space_after=14)
    
    future_work = [
        ("Concept Drift", "Scammers continuously adapt vocabularies, meaning models must be updated periodically with fresh training data."),
        ("Transformer Models (BERT)", "Moving from TF-IDF word frequencies to pre-trained transformer embeddings (BERT/RoBERTa) to capture deep semantic patterns."),
        ("Browser Extension Integration", "Deploying the trained model as a lightweight browser extension to flag fraudulent listings in real-time for users.")
    ]
    for title, desc in future_work:
        p_t = tf_c_r.add_paragraph()
        format_p(p_t, f"•  {title}", 14, bold=True, color=C_ACCENT, space_after=2)
        p_d = tf_c_r.add_paragraph()
        format_p(p_d, desc, 12, bold=False, color=C_TEXT, space_after=10)

    # ==================== SLIDE 11: Q&A PREPARATION ====================
    slide11 = prs.slides.add_slide(slide_layout)
    set_slide_background(slide11, C_BG)
    add_slide_header(slide11, "Q&A Session Preparation")

    # Split into two visual boxes
    q_left = slide11.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(5.6), Inches(4.8))
    tf_q_l = q_left.text_frame
    tf_q_l.word_wrap = True
    format_p(tf_q_l.paragraphs[0], "EXPECTED TECHNICAL QUESTIONS", 15, bold=True, color=C_PRIMARY, space_after=12)
    
    q_items1 = [
        ("Q: How did you address the class imbalance?", 
         "A: We used balanced class weights in Logistic Regression/Random Forest, and configured scale_pos_weight in XGBoost to penalize minority class errors heavily."),
        ("Q: Why not use Naive Bayes?", 
         "A: Naive Bayes assumes independent features, which fails when text is combined with metadata (like logo presence). XGBoost handles interactions better.")
    ]
    for q, a in q_items1:
        p_q = tf_q_l.add_paragraph()
        format_p(p_q, q, 13, bold=True, color=C_ACCENT, space_after=2)
        p_a = tf_q_l.add_paragraph()
        format_p(p_a, a, 12, bold=False, color=C_TEXT, space_after=12)

    q_right = slide11.shapes.add_textbox(Inches(6.8), Inches(1.8), Inches(5.7), Inches(4.8))
    tf_q_r = q_right.text_frame
    tf_q_r.word_wrap = True
    format_p(tf_q_r.paragraphs[0], "EXPECTED DOMAIN QUESTIONS", 15, bold=True, color=C_PRIMARY, space_after=12)
    
    q_items2 = [
        ("Q: Why does the presence of a company logo matter?", 
         "A: Legitimate employers invest time in setting up complete corporate profiles. Scammers or bots creating mass listings skip logo uploads."),
        ("Q: What is the risk of false positives?", 
         "A: XGBoost's high Precision (89.24%) guarantees that only ~10% of posts flagged as fake are actually real, ensuring legitimate posts aren't incorrectly blocked.")
    ]
    for q, a in q_items2:
        p_q = tf_q_r.add_paragraph()
        format_p(p_q, q, 13, bold=True, color=C_ACCENT, space_after=2)
        p_a = tf_q_r.add_paragraph()
        format_p(p_a, a, 12, bold=False, color=C_TEXT, space_after=12)

    # Save presentation
    output_dir = os.path.dirname(os.path.abspath(__file__)) if '__file__' in locals() else os.getcwd()
    output_path = os.path.join(output_dir, "presentation", "fake_job_posting_detection_presentation.pptx")
    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    
    prs.save(output_path)
    print(f"PowerPoint Presentation successfully created at: {output_path}")

if __name__ == "__main__":
    create_deck()
